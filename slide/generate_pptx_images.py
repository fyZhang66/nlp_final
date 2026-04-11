"""
Generate PPTX by screenshotting each HTML slide with Playwright.
Each slide becomes a full-bleed 2x-resolution image — pixel-perfect and non-editable.
"""

import os

SLIDE_DIR = os.path.dirname(os.path.abspath(__file__))
OUT_PPTX  = os.path.join(SLIDE_DIR, "ABSA_Final_Presentation.pptx")
IMG_DIR   = os.path.join(SLIDE_DIR, "_slide_images")
os.makedirs(IMG_DIR, exist_ok=True)

SLIDES = [f"slide{i}.html" for i in range(1, 17)]

# ── Step 1: Screenshot each HTML slide ───────────────────────────────
print("Step 1: Capturing 2x screenshots with Playwright …")

from playwright.sync_api import sync_playwright

with sync_playwright() as p:
    browser = p.chromium.launch(headless=True)
    # 2x device_scale_factor for crisp rendering in PPTX
    page = browser.new_page(
        viewport={"width": 1280, "height": 720},
        device_scale_factor=2,
    )

    for fname in SLIDES:
        html_path = os.path.join(SLIDE_DIR, fname)
        if not os.path.exists(html_path):
            print(f"  SKIP {fname} (not found)")
            continue
        file_url = f"file://{html_path}"
        page.goto(file_url, wait_until="networkidle")
        # Wait for Google Fonts and Font Awesome to fully load
        page.wait_for_timeout(1500)
        # Force all fonts loaded
        page.evaluate("() => document.fonts.ready")
        page.wait_for_timeout(300)

        out_png = os.path.join(IMG_DIR, fname.replace(".html", ".png"))
        page.screenshot(
            path=out_png,
            clip={"x": 0, "y": 0, "width": 1280, "height": 720},
        )
        print(f"  ✓ {fname}")

    browser.close()

print(f"\nAll screenshots saved to {IMG_DIR}")

# ── Step 2: Build PPTX from images ──────────────────────────────────
print("\nStep 2: Building image-based PPTX …")

from pptx import Presentation
from pptx.util import Inches

prs = Presentation()
prs.slide_width  = Inches(13.333)   # 16:9 widescreen
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]         # blank layout

for fname in SLIDES:
    png_path = os.path.join(IMG_DIR, fname.replace(".html", ".png"))
    if not os.path.exists(png_path):
        continue

    slide = prs.slides.add_slide(BLANK)
    slide.shapes.add_picture(
        png_path,
        left=0,
        top=0,
        width=prs.slide_width,
        height=prs.slide_height,
    )

prs.save(OUT_PPTX)
print(f"\n✓ Saved → {OUT_PPTX}  ({len(prs.slides)} slides, image-based / non-editable)")
