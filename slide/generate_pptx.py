"""Generate the final ABSA presentation as a .pptx file."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import json, os

# ── Colours ──────────────────────────────────────────────────────────────
BG       = RGBColor(0x18, 0x1B, 0x24)
PURPLE   = RGBColor(0xB1, 0x65, 0xFB)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
GRAY     = RGBColor(0xA0, 0xAE, 0xC0)
LGRAY    = RGBColor(0xE0, 0xE0, 0xE0)
GREEN    = RGBColor(0x4A, 0xDE, 0x80)
RED      = RGBColor(0xEF, 0x44, 0x44)
YELLOW   = RGBColor(0xFB, 0xBF, 0x24)
BLUE     = RGBColor(0x60, 0xA5, 0xFA)
LILAC    = RGBColor(0xC0, 0x84, 0xFC)
CARD_BG  = RGBColor(0x22, 0x25, 0x2F)
CARD_BD  = RGBColor(0x33, 0x36, 0x42)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

# Use blank layout
BLANK = prs.slide_layouts[6]

# ── Helper functions ─────────────────────────────────────────────────────

def set_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)

def set_text(tf, text, size=14, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font_name='Calibri'):
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return p

def add_para(tf, text, size=14, color=WHITE, bold=False, align=PP_ALIGN.LEFT, space_before=Pt(4), font_name='Calibri'):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    if space_before:
        p.space_before = space_before
    return p

def add_run(p, text, size=14, color=WHITE, bold=False, font_name='Calibri'):
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font_name
    return run

def purple_bar(slide, left, top, height=Inches(0.8)):
    return add_rect(slide, left, top, Inches(0.08), height, PURPLE)

def slide_title(slide, title_text, subtitle_text=None, left=Inches(0.9), top=Inches(0.55)):
    purple_bar(slide, left - Inches(0.18), top, Inches(0.85))
    tb = add_text_box(slide, left, top, Inches(10), Inches(0.55))
    set_text(tb.text_frame, title_text, size=36, bold=True, color=WHITE)
    if subtitle_text:
        tb2 = add_text_box(slide, left, top + Inches(0.55), Inches(10), Inches(0.35))
        set_text(tb2.text_frame, subtitle_text, size=16, color=GRAY)
    return top + Inches(1.1)

def footer(slide, text="NLP Final Project - INFO 7610"):
    tb = add_text_box(slide, Inches(10.5), Inches(7.0), Inches(2.5), Inches(0.35))
    set_text(tb.text_frame, text, size=11, color=GRAY, align=PP_ALIGN.RIGHT)


def make_table(slide, left, top, width, rows_data, col_widths, header_color=PURPLE):
    """rows_data: list of lists; first row is header."""
    n_rows = len(rows_data)
    n_cols = len(rows_data[0])
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, Inches(0.38 * n_rows))
    tbl = table_shape.table

    for ci, cw in enumerate(col_widths):
        tbl.columns[ci].width = cw

    for ri, row in enumerate(rows_data):
        for ci, cell_text in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.text = str(cell_text)
            cell.fill.solid()
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)
                p.font.name = 'Calibri'
                if ri == 0:
                    cell.fill.fore_color.rgb = RGBColor(0x25, 0x20, 0x3A)
                    p.font.color.rgb = PURPLE
                    p.font.bold = True
                else:
                    cell.fill.fore_color.rgb = BG if ri % 2 == 1 else RGBColor(0x1E, 0x21, 0x2B)
                    p.font.color.rgb = LGRAY
            cell.margin_left = Pt(6)
            cell.margin_right = Pt(6)
            cell.margin_top = Pt(4)
            cell.margin_bottom = Pt(4)

    # remove table borders (set to background)
    from pptx.oxml.ns import qn
    tbl_xml = tbl._tbl
    for tc in tbl_xml.iter(qn('a:tc')):
        tcPr = tc.find(qn('a:tcPr'))
        if tcPr is None:
            continue
        for edge in ['lnL', 'lnR', 'lnT', 'lnB']:
            ln = tcPr.find(qn(f'a:{edge}'))
            if ln is not None:
                ln.attrib['w'] = '0'

    return table_shape


def insight_box(slide, left, top, width, text, icon_char="*"):
    shape = add_rounded_rect(slide, left, top, width, Inches(0.65),
                             RGBColor(0x20, 0x1D, 0x30), PURPLE)
    tb = add_text_box(slide, left + Inches(0.15), top + Inches(0.08), width - Inches(0.3), Inches(0.55))
    set_text(tb.text_frame, text, size=12, color=LGRAY)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); set_bg(s)

# Purple accent bar
add_rect(s, Inches(0.72), Inches(1.4), Inches(0.08), Inches(2.2), PURPLE)

tb = add_text_box(s, Inches(0.95), Inches(0.7), Inches(5), Inches(0.4))
set_text(tb.text_frame, "NLP FINAL PROJECT", size=14, bold=True, color=PURPLE)

tb = add_text_box(s, Inches(0.95), Inches(1.4), Inches(9), Inches(0.6))
set_text(tb.text_frame, "End-to-End Aspect-Based", size=30, bold=True, color=PURPLE)

tb = add_text_box(s, Inches(0.95), Inches(2.0), Inches(10), Inches(0.7))
set_text(tb.text_frame, "Sentiment Analysis for Reviews", size=44, bold=True, color=WHITE)

tb = add_text_box(s, Inches(0.95), Inches(2.8), Inches(9), Inches(0.5))
set_text(tb.text_frame, "ATE + ASC Pipeline with Cross-Domain Evaluation on SemEval 2014", size=18, color=GRAY)

tb = add_text_box(s, Inches(0.95), Inches(3.5), Inches(9), Inches(0.4))
set_text(tb.text_frame, "INFO 7610 Natural Language Processing — Final Presentation  |  April 2026", size=13, color=GRAY)

# Team members
team = [
    ("Gousu Ding", "ding.go@northeastern.edu"),
    ("Fangyuan Zhang", "zhang.fangyua@northeastern.edu"),
    ("Yunzhu Chen", "chen.yunzh@northeastern.edu"),
    ("Yubo Wang", "wang.yubo3@northeastern.edu"),
]

tb = add_text_box(s, Inches(0.95), Inches(5.6), Inches(3), Inches(0.4))
set_text(tb.text_frame, "Team Members", size=16, bold=True, color=WHITE)

for i, (name, email) in enumerate(team):
    x = Inches(0.95 + i * 2.85)
    shape = add_rounded_rect(s, x, Inches(6.1), Inches(2.65), Inches(0.7), CARD_BG, RGBColor(0x45, 0x30, 0x70))
    tb = add_text_box(s, x + Inches(0.12), Inches(6.15), Inches(2.4), Inches(0.65))
    set_text(tb.text_frame, name, size=13, bold=True, color=WHITE)
    add_para(tb.text_frame, email, size=10, color=GRAY, space_before=Pt(2))

# Date
tb = add_text_box(s, Inches(10.5), Inches(0.7), Inches(2.5), Inches(0.6))
tf = tb.text_frame
set_text(tf, "Presentation Date", size=12, color=GRAY, align=PP_ALIGN.RIGHT)
add_para(tf, "April 2026", size=15, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)

footer(s)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 2 — Motivation
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); set_bg(s)
y = slide_title(s, "Why Star Ratings Are Not Enough", "ABSA reveals fine-grained, actionable insights from reviews")

# Review example
shape = add_rounded_rect(s, Inches(0.7), y, Inches(5.8), Inches(2.0), CARD_BG, PURPLE)
tb = add_text_box(s, Inches(0.9), y + Inches(0.12), Inches(5.4), Inches(1.8))
tf = tb.text_frame
set_text(tf, '"The pasta was absolutely delicious and the ambiance was lovely, but we waited 45 minutes for a table and the waiter was rude. 3 stars."', size=14, color=LGRAY)
add_para(tf, "", size=8, color=GRAY)
add_para(tf, "★★★☆☆  3.0", size=16, bold=True, color=YELLOW)

# ABSA extraction results
tb = add_text_box(s, Inches(7.0), y, Inches(5.5), Inches(0.4))
set_text(tb.text_frame, "ABSA Extraction Results", size=16, bold=True, color=PURPLE)

aspects = [
    ("pasta", "Positive", GREEN),
    ("ambiance", "Positive", GREEN),
    ("wait time", "Negative", RED),
    ("waiter", "Negative", RED),
]
for i, (asp, sent, clr) in enumerate(aspects):
    ay = y + Inches(0.5 + i * 0.45)
    tb = add_text_box(s, Inches(7.2), ay, Inches(2.5), Inches(0.4))
    set_text(tb.text_frame, asp, size=14, color=WHITE)
    tb2 = add_text_box(s, Inches(10.2), ay, Inches(2), Inches(0.4))
    set_text(tb2.text_frame, sent, size=14, bold=True, color=clr)

# Key insight
insight_box(s, Inches(0.7), Inches(5.5), Inches(11.9),
            "Key Insight: A single star rating hides actionable details. ABSA lets businesses see exactly what customers praise and complain about, enabling targeted improvements.")

footer(s)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 3 — Problem Definition
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); set_bg(s)
y = slide_title(s, "Problem Definition", "Two-Stage End-to-End ABSA Pipeline: ATE + ASC")

# Stage 1
shape = add_rounded_rect(s, Inches(0.7), y, Inches(5.8), Inches(2.3), CARD_BG, CARD_BD)
tb = add_text_box(s, Inches(0.9), y + Inches(0.1), Inches(5.5), Inches(2.1))
tf = tb.text_frame
p = set_text(tf, "", size=13, color=WHITE)
add_run(p, "1 ", size=16, bold=True, color=PURPLE)
add_run(p, "Aspect Term Extraction (ATE)", size=16, bold=True, color=WHITE)
add_para(tf, "Input: Raw review sentence", size=13, color=GRAY)
add_para(tf, "Task: BIO sequence labeling (like NER) to extract aspect spans", size=13, color=GRAY)
add_para(tf, "Output: Aspect terms with B-ASP / I-ASP / O tags", size=13, color=GRAY)
add_para(tf, "", size=6, color=GRAY)
add_para(tf, 'Example: "The [B-ASP]clam [I-ASP]chowder was incredible"', size=12, color=LGRAY)

# Stage 2
shape = add_rounded_rect(s, Inches(0.7), y + Inches(2.5), Inches(5.8), Inches(2.3), CARD_BG, CARD_BD)
tb = add_text_box(s, Inches(0.9), y + Inches(2.6), Inches(5.5), Inches(2.1))
tf = tb.text_frame
p = set_text(tf, "", size=13, color=WHITE)
add_run(p, "2 ", size=16, bold=True, color=PURPLE)
add_run(p, "Aspect Sentiment Classification (ASC)", size=16, bold=True, color=WHITE)
add_para(tf, "Input: (sentence, extracted aspect term)", size=13, color=GRAY)
add_para(tf, "Task: 3-class classification → Positive / Negative / Neutral", size=13, color=GRAY)
add_para(tf, 'Encoding: [CLS] sentence [SEP] aspect [SEP]', size=13, color=GRAY)
add_para(tf, "", size=6, color=GRAY)
add_para(tf, '("...clam chowder was incredible...", "clam chowder") → Positive', size=12, color=LGRAY)

# Pipeline flow on right
shape = add_rounded_rect(s, Inches(6.8), y, Inches(5.8), Inches(4.8), CARD_BG, CARD_BD)
tb = add_text_box(s, Inches(7.0), y + Inches(0.15), Inches(5.4), Inches(0.4))
set_text(tb.text_frame, "Pipeline Flow", size=16, bold=True, color=WHITE)

flow_items = ["Raw Text", "→", "ATE", "→", "Extracted\nAspects", "→", "ASC", "→", "(Aspect,\nSentiment)"]
for i, item in enumerate(flow_items):
    if item == "→":
        tb = add_text_box(s, Inches(7.0 + i * 0.58), y + Inches(1.0), Inches(0.5), Inches(0.5))
        set_text(tb.text_frame, "→", size=18, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)
    else:
        bx = Inches(6.85 + i * 0.58)
        is_model = item in ("ATE", "ASC")
        shape = add_rounded_rect(s, bx, y + Inches(0.8), Inches(0.55), Inches(0.65),
                                 PURPLE if is_model else RGBColor(0x2A, 0x2D, 0x38),
                                 PURPLE if is_model else CARD_BD)
        tb = add_text_box(s, bx + Inches(0.02), y + Inches(0.85), Inches(0.5), Inches(0.6))
        set_text(tb.text_frame, item, size=8, color=WHITE, align=PP_ALIGN.CENTER, bold=is_model)

tb = add_text_box(s, Inches(7.0), y + Inches(2.0), Inches(5.3), Inches(2.5))
tf = tb.text_frame
set_text(tf, "Key Points", size=14, bold=True, color=PURPLE)
add_para(tf, "• ATE extracts aspect terms via BIO labeling (same as NER)", size=12, color=LGRAY)
add_para(tf, "• ASC classifies sentiment using sentence-aspect pair encoding", size=12, color=LGRAY)
add_para(tf, "• Pipeline chains both stages end-to-end on raw text", size=12, color=LGRAY)
add_para(tf, "• No predefined aspect list needed", size=12, color=LGRAY)

footer(s)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 4 — Dataset
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); set_bg(s)
y = slide_title(s, "Dataset", "SemEval 2014 Task 4 — Standard ABSA Benchmark")

# Restaurant card
shape = add_rounded_rect(s, Inches(0.7), y, Inches(5.8), Inches(1.5), CARD_BG, CARD_BD)
tb = add_text_box(s, Inches(0.9), y + Inches(0.1), Inches(5.5), Inches(1.3))
tf = tb.text_frame
set_text(tf, "Restaurant Domain", size=16, bold=True, color=WHITE)
add_para(tf, "~3,000 train / ~800 test sentences  |  1,120 test aspect-sentiment pairs", size=12, color=GRAY)
add_para(tf, "Aspects: food, service, ambiance, price", size=12, color=GRAY)
add_para(tf, '"The pasta is out of this world, but prices are steep."', size=11, color=RGBColor(0x80, 0x88, 0x98))

# Laptop card
shape = add_rounded_rect(s, Inches(0.7), y + Inches(1.7), Inches(5.8), Inches(1.5), CARD_BG, CARD_BD)
tb = add_text_box(s, Inches(0.9), y + Inches(1.8), Inches(5.5), Inches(1.3))
tf = tb.text_frame
set_text(tf, "Laptop Domain", size=16, bold=True, color=WHITE)
add_para(tf, "~3,000 train / ~800 test sentences  |  638 test aspect-sentiment pairs", size=12, color=GRAY)
add_para(tf, "Aspects: screen, battery, performance, design", size=12, color=GRAY)
add_para(tf, '"Battery life is great but the keyboard feels cheap."', size=11, color=RGBColor(0x80, 0x88, 0x98))

# Label formats
shape = add_rounded_rect(s, Inches(6.8), y, Inches(5.8), Inches(1.5), CARD_BG, RGBColor(0x45, 0x30, 0x70))
tb = add_text_box(s, Inches(7.0), y + Inches(0.1), Inches(5.4), Inches(1.3))
tf = tb.text_frame
set_text(tf, "Label Formats", size=14, bold=True, color=PURPLE)
add_para(tf, "ATE:  BIO tagging — B-ASP / I-ASP / O", size=12, color=LGRAY)
add_para(tf, "ASC:  3-class — Positive / Negative / Neutral", size=12, color=LGRAY)
add_para(tf, "Train/Val split: 90/10, Seed = 42  |  Conflicts skipped", size=11, color=GRAY)

# Distribution
shape = add_rounded_rect(s, Inches(6.8), y + Inches(1.7), Inches(5.8), Inches(1.5), CARD_BG, RGBColor(0x45, 0x30, 0x70))
tb = add_text_box(s, Inches(7.0), y + Inches(1.8), Inches(5.4), Inches(1.3))
tf = tb.text_frame
set_text(tf, "ASC Label Distribution", size=14, bold=True, color=PURPLE)
add_para(tf, "Positive:  ~46%  ██████████████████░░░░░░░░░░░░░░░░░░░░", size=11, color=GREEN)
add_para(tf, "Negative:  ~34%  █████████████░░░░░░░░░░░░░░░░░░░░░░░░░", size=11, color=RED)
add_para(tf, "Neutral:   ~20%  ████████░░░░░░░░░░░░░░░░░░░░░░░░░░░░░░", size=11, color=YELLOW)

insight_box(s, Inches(0.7), Inches(5.6), Inches(11.9),
            "Class imbalance: Neutral is underrepresented (~20%), which may bias ASC predictions toward Positive/Negative.")
footer(s)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 5 — ATE Architecture
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); set_bg(s)
y = slide_title(s, "Model Architecture", "Stage 1: Aspect Term Extraction (ATE) — Token Classification")

# Architecture box
shape = add_rounded_rect(s, Inches(0.7), y, Inches(6.2), Inches(2.0), CARD_BG, CARD_BD)
tb = add_text_box(s, Inches(0.9), y + Inches(0.1), Inches(5.8), Inches(1.8))
tf = tb.text_frame
set_text(tf, "Architecture", size=14, bold=True, color=PURPLE)
add_para(tf, "Backbone: BERT-base-uncased or DeBERTa-base", size=12, color=LGRAY)
add_para(tf, "Head: Linear layer over each token → {B-ASP, I-ASP, O}", size=12, color=LGRAY)
add_para(tf, "Subword alignment: first subword inherits word label; continuations = -100", size=12, color=LGRAY)
add_para(tf, "Evaluation: entity-level F1 via seqeval (exact span match)", size=12, color=LGRAY)

# Flow diagram
flow_labels = ["Input Tokens", "BERT/DeBERTa", "Token Embs", "Linear Head", "BIO Labels"]
for i, lbl in enumerate(flow_labels):
    bx = Inches(0.7 + i * 1.35)
    is_model = i in (1, 3)
    shape = add_rounded_rect(s, bx, y + Inches(2.2), Inches(1.15), Inches(0.5),
                             PURPLE if is_model else RGBColor(0x2A, 0x2D, 0x38),
                             PURPLE if is_model else CARD_BD)
    tb = add_text_box(s, bx, y + Inches(2.25), Inches(1.15), Inches(0.45))
    set_text(tb.text_frame, lbl, size=10, color=WHITE, bold=is_model, align=PP_ALIGN.CENTER)
    if i < len(flow_labels) - 1:
        tb = add_text_box(s, bx + Inches(1.15), y + Inches(2.3), Inches(0.2), Inches(0.35))
        set_text(tb.text_frame, "→", size=14, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)

# Example
shape = add_rounded_rect(s, Inches(0.7), y + Inches(2.9), Inches(6.2), Inches(1.0), RGBColor(0x20, 0x1D, 0x30), PURPLE)
tb = add_text_box(s, Inches(0.9), y + Inches(2.95), Inches(5.8), Inches(0.9))
tf = tb.text_frame
set_text(tf, "Example BIO Tagging", size=11, bold=True, color=PURPLE)
add_para(tf, 'The [B-ASP]clam [I-ASP]chowder was incredible but the [B-ASP]waiter was rude .', size=11, color=LGRAY)
add_para(tf, 'Extracted:  "clam chowder",  "waiter"', size=11, color=GREEN)

# Hyperparameters table
hp_data = [
    ["Parameter", "BERT", "DeBERTa"],
    ["Pre-trained Model", "bert-base-uncased", "microsoft/deberta-base"],
    ["Learning Rate", "3e-5", "2e-5"],
    ["Epochs", "5", "5"],
    ["Train Batch Size", "16", "8"],
    ["Eval Batch Size", "32", "16"],
    ["Max Seq Length", "128", "128"],
    ["Metric for Best", "eval_f1", "eval_f1"],
]
make_table(s, Inches(7.2), y, Inches(5.4), hp_data,
           [Inches(2.0), Inches(1.7), Inches(1.7)])

insight_box(s, Inches(7.2), y + Inches(3.3), Inches(5.4),
            "Uses HuggingFace AutoModelForTokenClassification with Trainer API. DataCollator handles dynamic padding.")
footer(s)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 6 — ASC Architecture
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); set_bg(s)
y = slide_title(s, "Model Architecture", "Stage 2: Aspect Sentiment Classification (ASC) — Sequence Classification")

# Architecture box
shape = add_rounded_rect(s, Inches(0.7), y, Inches(6.2), Inches(1.8), CARD_BG, CARD_BD)
tb = add_text_box(s, Inches(0.9), y + Inches(0.1), Inches(5.8), Inches(1.6))
tf = tb.text_frame
set_text(tf, "Architecture", size=14, bold=True, color=PURPLE)
add_para(tf, "Input: Sentence-aspect pair encoding", size=12, color=LGRAY)
add_para(tf, "Backbone: BERT-base-uncased or DeBERTa-base", size=12, color=LGRAY)
add_para(tf, "Head: Linear 3-class classifier on [CLS] representation", size=12, color=LGRAY)
add_para(tf, "Evaluation: Accuracy + Macro-F1 (balanced across classes)", size=12, color=LGRAY)

# Input encoding diagram
enc_labels = ["[CLS]", "review sentence tokens", "[SEP]", "aspect term", "[SEP]"]
enc_widths = [0.65, 2.8, 0.65, 1.4, 0.65]
enc_colors = [PURPLE, RGBColor(0x1A, 0x3A, 0x25), PURPLE, RGBColor(0x3A, 0x2E, 0x15), PURPLE]
ex = Inches(0.7)
for i, (lbl, w, clr) in enumerate(zip(enc_labels, enc_widths, enc_colors)):
    shape = add_rounded_rect(s, ex, y + Inches(2.0), Inches(w), Inches(0.5), clr,
                             PURPLE if clr == PURPLE else CARD_BD)
    tb = add_text_box(s, ex, y + Inches(2.05), Inches(w), Inches(0.45))
    set_text(tb.text_frame, lbl, size=10, color=WHITE, bold=(clr == PURPLE), align=PP_ALIGN.CENTER)
    ex += Inches(w + 0.03)

# Example
shape = add_rounded_rect(s, Inches(0.7), y + Inches(2.7), Inches(6.2), Inches(1.2), RGBColor(0x20, 0x1D, 0x30), PURPLE)
tb = add_text_box(s, Inches(0.9), y + Inches(2.75), Inches(5.8), Inches(1.1))
tf = tb.text_frame
set_text(tf, "Examples", size=11, bold=True, color=PURPLE)
add_para(tf, '("The clam chowder was incredible but the waiter was rude.", "clam chowder") → Positive', size=11, color=LGRAY)
add_para(tf, '("The clam chowder was incredible but the waiter was rude.", "waiter") → Negative', size=11, color=LGRAY)

# Hyperparameters table
hp_data = [
    ["Parameter", "BERT", "DeBERTa"],
    ["Pre-trained Model", "bert-base-uncased", "microsoft/deberta-base"],
    ["Learning Rate", "3e-5", "2e-5"],
    ["Epochs", "5", "3"],
    ["Batch Size (train)", "16", "16"],
    ["Batch Size (eval)", "32", "32"],
    ["Max Seq Length", "128", "128"],
    ["Weight Decay", "0.01", "0.01"],
    ["Warmup Steps", "100", "100"],
    ["Metric for Best", "eval_macro_f1", "eval_macro_f1"],
]
make_table(s, Inches(7.2), y, Inches(5.4), hp_data,
           [Inches(2.0), Inches(1.7), Inches(1.7)])

insight_box(s, Inches(7.2), y + Inches(3.6), Inches(5.4),
            "Uses HuggingFace AutoModelForSequenceClassification. Sentence-pair tokenization: tokenizer(sentence, aspect)")
footer(s)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 7 — In-Domain ATE Results
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); set_bg(s)
y = slide_title(s, "In-Domain ATE Results", "Standalone Aspect Term Extraction — Entity-level F1 (seqeval, exact span match)")

ate_data = [
    ["Domain", "Model", "Precision", "Recall", "F1 Score"],
    ["Restaurant", "BERT", "81.2%", "85.5%", "83.3%"],
    ["Restaurant", "DeBERTa", "84.0%", "88.4%", "86.2%"],
    ["Laptop", "BERT", "74.5%", "81.0%", "77.6%"],
    ["Laptop", "DeBERTa", "80.5%", "83.1%", "81.8%"],
]
make_table(s, Inches(0.7), y, Inches(11.9), ate_data,
           [Inches(2.0), Inches(2.5), Inches(2.0), Inches(2.0), Inches(2.5)])

# Takeaway cards
cards = [
    ("DeBERTa Wins", "Outperforms BERT by +2.9% (rest.) and +4.2% (laptop) in F1. Disentangled attention helps with aspect boundary detection."),
    ("Restaurant Easier", "Higher F1 across both models. Restaurant aspects (food, service) have more consistent patterns than laptop aspects."),
    ("Recall > Precision", "Models tend to over-extract (higher recall), generating some false positive aspect terms."),
]
for i, (title, text) in enumerate(cards):
    cx = Inches(0.7 + i * 4.1)
    shape = add_rounded_rect(s, cx, y + Inches(2.8), Inches(3.8), Inches(1.5), CARD_BG, CARD_BD)
    tb = add_text_box(s, cx + Inches(0.15), y + Inches(2.9), Inches(3.5), Inches(1.3))
    tf = tb.text_frame
    set_text(tf, title, size=14, bold=True, color=PURPLE)
    add_para(tf, text, size=11, color=GRAY)

footer(s)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 8 — In-Domain ASC Results
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); set_bg(s)
y = slide_title(s, "In-Domain ASC Results", "Standalone Aspect Sentiment Classification — Gold Aspect Terms")

asc_data = [
    ["Domain", "Model", "Accuracy", "Macro-F1", "Pos F1", "Neg F1", "Neu F1"],
    ["Restaurant", "BERT", "83.0%", "74.4%", "0.91", "0.76", "0.56"],
    ["Restaurant", "DeBERTa", "84.2%", "74.8%", "0.91", "0.80", "0.53"],
    ["Laptop", "BERT", "75.7%", "69.7%", "0.88", "0.71", "0.51"],
    ["Laptop", "DeBERTa", "79.5%", "75.5%", "0.89", "0.73", "0.64"],
]
make_table(s, Inches(0.7), y, Inches(11.9), asc_data,
           [Inches(1.8), Inches(1.8), Inches(1.5), Inches(1.5), Inches(1.5), Inches(1.5), Inches(1.5)])

# Per-class bars
shape = add_rounded_rect(s, Inches(0.7), y + Inches(2.6), Inches(5.6), Inches(2.2), CARD_BG, CARD_BD)
tb = add_text_box(s, Inches(0.9), y + Inches(2.7), Inches(5.2), Inches(2.0))
tf = tb.text_frame
set_text(tf, "Per-Class F1 (Restaurant DeBERTa)", size=13, bold=True, color=PURPLE)
add_para(tf, "Positive:   0.91  ████████████████████████████████████░░░░", size=11, color=GREEN)
add_para(tf, "Negative:   0.80  ████████████████████████████░░░░░░░░░░░░", size=11, color=RED)
add_para(tf, "Neutral:    0.53  ██████████████████░░░░░░░░░░░░░░░░░░░░░░", size=11, color=YELLOW)

shape = add_rounded_rect(s, Inches(6.6), y + Inches(2.6), Inches(5.6), Inches(2.2), CARD_BG, CARD_BD)
tb = add_text_box(s, Inches(6.8), y + Inches(2.7), Inches(5.2), Inches(2.0))
tf = tb.text_frame
set_text(tf, "Per-Class F1 (Laptop DeBERTa)", size=13, bold=True, color=PURPLE)
add_para(tf, "Positive:   0.89  ███████████████████████████████████░░░░░", size=11, color=GREEN)
add_para(tf, "Negative:   0.73  █████████████████████████░░░░░░░░░░░░░░░", size=11, color=RED)
add_para(tf, "Neutral:    0.64  ██████████████████████░░░░░░░░░░░░░░░░░░", size=11, color=YELLOW)

insight_box(s, Inches(0.7), y + Inches(5.0), Inches(11.9),
            "Key insight: Neutral is the hardest class (F1: 0.51-0.64). Only ~20% of training data is Neutral. DeBERTa improves Laptop Neutral by +13 points (0.51 → 0.64).")
footer(s)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 9 — Cross-Domain Experiment Matrix
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); set_bg(s)
y = slide_title(s, "Cross-Domain Experiments", "8-Run Evaluation Matrix: 2 Domains × 2 Test Sets × 2 Backbones")

matrix_data = [
    ["ID", "Train Domain", "Test Domain", "Model", "Type"],
    ["1", "Restaurant", "Restaurant", "BERT", "In-domain"],
    ["2", "Restaurant", "Restaurant", "DeBERTa", "In-domain"],
    ["3", "Restaurant", "Laptop", "BERT", "Cross-domain"],
    ["4", "Restaurant", "Laptop", "DeBERTa", "Cross-domain"],
    ["5", "Laptop", "Restaurant", "BERT", "Cross-domain"],
    ["6", "Laptop", "Restaurant", "DeBERTa", "Cross-domain"],
    ["7", "Laptop", "Laptop", "BERT", "In-domain"],
    ["8", "Laptop", "Laptop", "DeBERTa", "In-domain"],
]
make_table(s, Inches(0.7), y, Inches(7.5), matrix_data,
           [Inches(0.6), Inches(1.8), Inches(1.8), Inches(1.5), Inches(1.8)])

# Visual 2x2 matrix
shape = add_rounded_rect(s, Inches(8.5), y, Inches(4.2), Inches(3.0), CARD_BG, CARD_BD)
tb = add_text_box(s, Inches(8.7), y + Inches(0.1), Inches(3.8), Inches(0.4))
set_text(tb.text_frame, "Domain Transfer Matrix", size=14, bold=True, color=PURPLE)

# Headers
tb = add_text_box(s, Inches(9.6), y + Inches(0.5), Inches(1.3), Inches(0.35))
set_text(tb.text_frame, "Test: Rest.", size=10, color=GRAY, align=PP_ALIGN.CENTER)
tb = add_text_box(s, Inches(10.9), y + Inches(0.5), Inches(1.3), Inches(0.35))
set_text(tb.text_frame, "Test: Lapt.", size=10, color=GRAY, align=PP_ALIGN.CENTER)

# Row labels + cells
tb = add_text_box(s, Inches(8.6), y + Inches(1.0), Inches(1.0), Inches(0.7))
set_text(tb.text_frame, "Train:\nRest.", size=10, color=GRAY, align=PP_ALIGN.RIGHT)
tb = add_text_box(s, Inches(8.6), y + Inches(1.8), Inches(1.0), Inches(0.7))
set_text(tb.text_frame, "Train:\nLapt.", size=10, color=GRAY, align=PP_ALIGN.RIGHT)

# In-domain cells (green tint)
shape = add_rounded_rect(s, Inches(9.7), y + Inches(1.0), Inches(1.1), Inches(0.7),
                         RGBColor(0x1A, 0x2A, 0x1A), RGBColor(0x2A, 0x5A, 0x2A))
tb = add_text_box(s, Inches(9.7), y + Inches(1.05), Inches(1.1), Inches(0.65))
tf = tb.text_frame
set_text(tf, "In-domain", size=10, color=GREEN, align=PP_ALIGN.CENTER, bold=True)
add_para(tf, "#1, #2", size=9, color=GRAY, align=PP_ALIGN.CENTER)

shape = add_rounded_rect(s, Inches(11.0), y + Inches(1.8), Inches(1.1), Inches(0.7),
                         RGBColor(0x1A, 0x2A, 0x1A), RGBColor(0x2A, 0x5A, 0x2A))
tb = add_text_box(s, Inches(11.0), y + Inches(1.85), Inches(1.1), Inches(0.65))
tf = tb.text_frame
set_text(tf, "In-domain", size=10, color=GREEN, align=PP_ALIGN.CENTER, bold=True)
add_para(tf, "#7, #8", size=9, color=GRAY, align=PP_ALIGN.CENTER)

# Cross-domain cells (yellow tint)
shape = add_rounded_rect(s, Inches(11.0), y + Inches(1.0), Inches(1.1), Inches(0.7),
                         RGBColor(0x2A, 0x25, 0x15), RGBColor(0x5A, 0x4A, 0x1A))
tb = add_text_box(s, Inches(11.0), y + Inches(1.05), Inches(1.1), Inches(0.65))
tf = tb.text_frame
set_text(tf, "Cross", size=10, color=YELLOW, align=PP_ALIGN.CENTER, bold=True)
add_para(tf, "#3, #4", size=9, color=GRAY, align=PP_ALIGN.CENTER)

shape = add_rounded_rect(s, Inches(9.7), y + Inches(1.8), Inches(1.1), Inches(0.7),
                         RGBColor(0x2A, 0x25, 0x15), RGBColor(0x5A, 0x4A, 0x1A))
tb = add_text_box(s, Inches(9.7), y + Inches(1.85), Inches(1.1), Inches(0.65))
tf = tb.text_frame
set_text(tf, "Cross", size=10, color=YELLOW, align=PP_ALIGN.CENTER, bold=True)
add_para(tf, "#5, #6", size=9, color=GRAY, align=PP_ALIGN.CENTER)

# Research question
insight_box(s, Inches(8.5), y + Inches(3.3), Inches(4.2),
            "Research Question: How well do ABSA models generalize across product domains? Where does the pipeline break?")

footer(s)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 10 — Cross-Domain Pipeline Results
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); set_bg(s)
y = slide_title(s, "Cross-Domain Pipeline Results", "End-to-End ATE → ASC Evaluation Across All 8 Experiments")

results_data = [
    ["Setting", "Model", "Type", "ATE F1", "ASC Acc (gold)", "ASC Acc (pred)", "Error Gap"],
    ["Rest → Rest", "BERT", "In", "56.1%", "83.0%", "84.8%", "-1.8%"],
    ["Rest → Rest", "DeBERTa", "In", "57.6%", "84.2%", "84.1%", "+0.2%"],
    ["Rest → Lapt", "BERT", "Cross", "25.5%", "76.2%", "78.5%", "-2.3%"],
    ["Rest → Lapt", "DeBERTa", "Cross", "32.5%", "79.5%", "82.0%", "-2.5%"],
    ["Lapt → Rest", "BERT", "Cross", "25.8%", "80.2%", "92.2%", "-12.0%"],
    ["Lapt → Rest", "DeBERTa", "Cross", "36.1%", "79.8%", "88.4%", "-8.6%"],
    ["Lapt → Lapt", "BERT", "In", "52.3%", "75.7%", "77.6%", "-1.8%"],
    ["Lapt → Lapt", "DeBERTa", "In", "54.9%", "79.5%", "80.1%", "-0.6%"],
]
make_table(s, Inches(0.7), y, Inches(11.9), results_data,
           [Inches(1.7), Inches(1.3), Inches(0.8), Inches(1.5), Inches(2.0), Inches(2.0), Inches(1.7)])

cards_text = [
    "ATE drops dramatically cross-domain: 25-36% F1 vs 52-58% in-domain. Aspect vocabularies are domain-specific.",
    "Anomaly: Lapt→Rest BERT achieves 92.2% ASC on predicted aspects — higher than gold! ATE extracts only high-confidence easy aspects.",
    "Note: Pipeline ATE F1 uses term-level exact string match (different from standalone seqeval entity-level F1).",
]
for i, text in enumerate(cards_text):
    cx = Inches(0.7 + i * 4.1)
    shape = add_rounded_rect(s, cx, y + Inches(3.6), Inches(3.8), Inches(1.1), CARD_BG, CARD_BD)
    tb = add_text_box(s, cx + Inches(0.12), y + Inches(3.65), Inches(3.55), Inches(1.0))
    set_text(tb.text_frame, text, size=11, color=LGRAY)

footer(s)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 11 — Key Finding
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); set_bg(s)
y = slide_title(s, "Key Finding", "ATE Degrades Cross-Domain, but ASC Remains Robust")

# ATE card (red)
shape = add_rounded_rect(s, Inches(0.7), y, Inches(5.8), Inches(1.8),
                         RGBColor(0x25, 0x18, 0x18), RGBColor(0x5A, 0x22, 0x22))
tb = add_text_box(s, Inches(0.9), y + Inches(0.1), Inches(5.4), Inches(1.6))
tf = tb.text_frame
set_text(tf, "ATE F1 — Degrades Significantly", size=16, bold=True, color=RED)
add_para(tf, "In-domain range:          52 – 58%", size=13, color=LGRAY)
add_para(tf, "Cross-domain range:     25 – 36%    (up to 56% relative drop)", size=13, color=RED, bold=True)

# ASC card (green)
shape = add_rounded_rect(s, Inches(0.7), y + Inches(2.0), Inches(5.8), Inches(1.8),
                         RGBColor(0x15, 0x25, 0x18), RGBColor(0x22, 0x5A, 0x2A))
tb = add_text_box(s, Inches(0.9), y + Inches(2.1), Inches(5.4), Inches(1.6))
tf = tb.text_frame
set_text(tf, "ASC Accuracy — Surprisingly Stable", size=16, bold=True, color=GREEN)
add_para(tf, "In-domain range:          76 – 85%", size=13, color=LGRAY)
add_para(tf, "Cross-domain range:     78 – 92%    (remains high!)", size=13, color=GREEN, bold=True)

# Why box
shape = add_rounded_rect(s, Inches(6.8), y, Inches(5.8), Inches(2.8), CARD_BG, CARD_BD)
tb = add_text_box(s, Inches(7.0), y + Inches(0.1), Inches(5.4), Inches(2.6))
tf = tb.text_frame
set_text(tf, "Why?", size=16, bold=True, color=PURPLE)
add_para(tf, "", size=6, color=WHITE)
add_para(tf, "✓  ASC transfers well: [CLS] sentence [SEP] aspect [SEP] captures sentiment patterns that are universal. Words like \"great\", \"terrible\", \"okay\" carry the same sentiment across domains.", size=12, color=LGRAY)
add_para(tf, "", size=4, color=WHITE)
add_para(tf, "✗  ATE is domain-specific: \"battery\" is an aspect in laptops but not restaurants. Aspect vocabularies don't transfer across product domains.", size=12, color=LGRAY)
add_para(tf, "", size=4, color=WHITE)
add_para(tf, "→  Implication: For cross-domain ABSA, invest in better ATE or domain-specific training data. ASC can reuse existing models.", size=12, color=LGRAY)

# Anomaly
shape = add_rounded_rect(s, Inches(6.8), y + Inches(3.0), Inches(5.8), Inches(0.8),
                         RGBColor(0x2A, 0x25, 0x15), RGBColor(0x5A, 0x4A, 0x1A))
tb = add_text_box(s, Inches(7.0), y + Inches(3.05), Inches(5.4), Inches(0.7))
set_text(tb.text_frame,
         "Anomaly: Lapt→Rest BERT gets 92.2% ASC accuracy on predicted aspects (higher than gold 80.2%!). ATE extracts only high-confidence easy aspects, giving ASC an easier subset.",
         size=11, color=LGRAY)

footer(s)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 12 — ATE Error Analysis
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); set_bg(s)
y = slide_title(s, "Error Analysis: ATE", "Categorizing Aspect Term Extraction Errors")

errors = [
    ("Correct  (~50-55%)", "Exact span match: Gold \"pasta\" → Pred \"pasta\"", GREEN),
    ("Missing  (~6-12%)", "Gold aspect not predicted: Gold \"wait time\" → Pred (nothing)", RED),
    ("Spurious  (~8-15%)", "Predicted aspect not in gold: Gold (none) → Pred \"restaurant\"", YELLOW),
    ("Boundary  (~34-40%)", "Partial overlap: Gold \"clam chowder\" → Pred \"chowder\" (partial)", BLUE),
]
for i, (name, desc, clr) in enumerate(errors):
    cy = y + i * Inches(1.0)
    shape = add_rounded_rect(s, Inches(0.7), cy, Inches(6.0), Inches(0.85), CARD_BG, clr)
    tb = add_text_box(s, Inches(0.9), cy + Inches(0.05), Inches(5.6), Inches(0.75))
    tf = tb.text_frame
    set_text(tf, name, size=14, bold=True, color=clr)
    add_para(tf, desc, size=11, color=GRAY)

# Distribution bars
shape = add_rounded_rect(s, Inches(7.0), y, Inches(5.6), Inches(2.5), CARD_BG, CARD_BD)
tb = add_text_box(s, Inches(7.2), y + Inches(0.1), Inches(5.2), Inches(2.3))
tf = tb.text_frame
set_text(tf, "Error Distribution (In-Domain)", size=14, bold=True, color=PURPLE)
add_para(tf, "", size=6, color=WHITE)
add_para(tf, "Correct     ████████████████████████████████░░░░░░░░░░░░░░  52%", size=11, color=GREEN)
add_para(tf, "Boundary   ███████████████████████░░░░░░░░░░░░░░░░░░░░░░░  37%", size=11, color=BLUE)
add_para(tf, "Spurious    ███████░░░░░░░░░░░░░░░░░░░░░░░░░░░░░░░░░░░░░░  12%", size=11, color=YELLOW)
add_para(tf, "Missing     █████░░░░░░░░░░░░░░░░░░░░░░░░░░░░░░░░░░░░░░░░░   8%", size=11, color=RED)

insight_box(s, Inches(7.0), y + Inches(2.7), Inches(5.6),
            "Key insight: Boundary errors are the dominant failure mode (~37%). Multi-word aspects like \"clam chowder\" are often partially extracted. This is the biggest opportunity for improvement.")
footer(s)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 13 — E2E Error Tracing
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); set_bg(s)
y = slide_title(s, "End-to-End Error Tracing", "How ATE Errors Propagate Through the Pipeline")

traces = [
    ("ATE Miss → Sentiment Lost", "20-30%", "Gold aspect not extracted. Sentiment completely missed.", RED),
    ("ATE Boundary → Wrong Sentiment", "20-30%", "Partial extraction leads to wrong sentiment.", BLUE),
    ("ATE Spurious → Extra Prediction", "10-20%", "False positive aspect generates phantom sentiment.", YELLOW),
    ("Correct ATE → Sentiment Mismatch", "5-15%", "Aspect correct but ASC predicts wrong polarity.", LILAC),
]
for i, (name, pct, desc, clr) in enumerate(traces):
    cy = y + i * Inches(0.95)
    shape = add_rounded_rect(s, Inches(0.7), cy, Inches(6.2), Inches(0.8), CARD_BG, clr)
    tb = add_text_box(s, Inches(0.9), cy + Inches(0.05), Inches(4.5), Inches(0.7))
    tf = tb.text_frame
    set_text(tf, name, size=13, bold=True, color=clr)
    add_para(tf, desc, size=11, color=GRAY)
    tb2 = add_text_box(s, Inches(5.5), cy + Inches(0.15), Inches(1.2), Inches(0.4))
    set_text(tb2.text_frame, pct, size=14, bold=True, color=clr, align=PP_ALIGN.RIGHT)

# Error source
shape = add_rounded_rect(s, Inches(7.2), y, Inches(5.4), Inches(1.5), CARD_BG, CARD_BD)
tb = add_text_box(s, Inches(7.4), y + Inches(0.1), Inches(5.0), Inches(1.3))
tf = tb.text_frame
set_text(tf, "Error Source Distribution", size=14, bold=True, color=PURPLE)
add_para(tf, "", size=6, color=WHITE)
p = add_para(tf, "", size=13, color=RED, bold=True)
add_run(p, "ATE errors:  ~85% of all E2E errors", size=13, color=RED, bold=True)
p = add_para(tf, "", size=13, color=GREEN, bold=True)
add_run(p, "ASC errors:  ~15% of all E2E errors", size=13, color=GREEN, bold=True)

# Distribution bars
shape = add_rounded_rect(s, Inches(7.2), y + Inches(1.7), Inches(5.4), Inches(1.8), CARD_BG, CARD_BD)
tb = add_text_box(s, Inches(7.4), y + Inches(1.8), Inches(5.0), Inches(1.6))
tf = tb.text_frame
set_text(tf, "Error Breakdown", size=13, bold=True, color=PURPLE)
add_para(tf, "ATE Miss → Lost         ███████████████░░░░░░░░░░░░░░░░  28%", size=10, color=RED)
add_para(tf, "ATE Boundary → Wrong  ████████████░░░░░░░░░░░░░░░░░░░░  25%", size=10, color=BLUE)
add_para(tf, "ATE Spurious → Extra   █████████░░░░░░░░░░░░░░░░░░░░░░░  18%", size=10, color=YELLOW)
add_para(tf, "Correct → Mismatch     ██████░░░░░░░░░░░░░░░░░░░░░░░░░░  12%", size=10, color=LILAC)

insight_box(s, Inches(7.2), y + Inches(3.7), Inches(5.4),
            "Most E2E errors (~85%) originate from ATE, not ASC. Improving extraction quality would have the biggest impact on overall performance.")

footer(s)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 14 — BERT vs DeBERTa
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); set_bg(s)
y = slide_title(s, "BERT vs DeBERTa", "Head-to-Head Comparison Across All Experiments")

# Model cards side by side
shape = add_rounded_rect(s, Inches(0.7), y, Inches(2.8), Inches(1.2), CARD_BG, RGBColor(0x30, 0x40, 0x60))
tb = add_text_box(s, Inches(0.9), y + Inches(0.1), Inches(2.4), Inches(1.0))
tf = tb.text_frame
set_text(tf, "BERT", size=22, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
add_para(tf, "bert-base-uncased\n110M parameters", size=11, color=GRAY, align=PP_ALIGN.CENTER)

tb = add_text_box(s, Inches(3.6), y + Inches(0.3), Inches(0.5), Inches(0.5))
set_text(tb.text_frame, "vs", size=18, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)

shape = add_rounded_rect(s, Inches(4.2), y, Inches(2.8), Inches(1.2), CARD_BG, RGBColor(0x40, 0x28, 0x60))
tb = add_text_box(s, Inches(4.4), y + Inches(0.1), Inches(2.4), Inches(1.0))
tf = tb.text_frame
set_text(tf, "DeBERTa", size=22, bold=True, color=LILAC, align=PP_ALIGN.CENTER)
add_para(tf, "microsoft/deberta-base\n139M parameters", size=11, color=GRAY, align=PP_ALIGN.CENTER)

# Gain table
gain_data = [
    ["Setting", "ATE F1 Gain", "ASC Acc Gain"],
    ["Rest → Rest", "+1.5 pp", "+1.2 pp"],
    ["Rest → Lapt", "+7.0 pp", "+3.5 pp"],
    ["Lapt → Rest", "+10.3 pp", "-3.8 pp *"],
    ["Lapt → Lapt", "+2.5 pp", "+2.6 pp"],
]
make_table(s, Inches(0.7), y + Inches(1.5), Inches(6.3), gain_data,
           [Inches(2.1), Inches(2.1), Inches(2.1)])

tb = add_text_box(s, Inches(0.7), y + Inches(3.6), Inches(6.3), Inches(0.35))
set_text(tb.text_frame, "* BERT's anomalously high 92.2% from easy-subset selection effect", size=10, color=GRAY)

# Why DeBERTa wins
shape = add_rounded_rect(s, Inches(7.2), y, Inches(5.4), Inches(3.0), CARD_BG, CARD_BD)
tb = add_text_box(s, Inches(7.4), y + Inches(0.1), Inches(5.0), Inches(2.8))
tf = tb.text_frame
set_text(tf, "Why DeBERTa Wins", size=16, bold=True, color=PURPLE)
add_para(tf, "", size=6, color=WHITE)
add_para(tf, "Disentangled attention separates content and position signals, helping distinguish tokens in multi-aspect sentences more precisely.", size=12, color=LGRAY)
add_para(tf, "", size=6, color=WHITE)
add_para(tf, "Better boundary detection for multi-word aspects. Relative position encoding helps identify where aspect spans begin and end.", size=12, color=LGRAY)
add_para(tf, "", size=6, color=WHITE)
add_para(tf, "More robust to domain shift in ATE. Biggest gain: Lapt→Rest ATE F1 jumps from 25.8% to 36.1% (+10.3 percentage points).", size=12, color=LGRAY)

# Tradeoff
shape = add_rounded_rect(s, Inches(7.2), y + Inches(3.2), Inches(5.4), Inches(0.7),
                         RGBColor(0x2A, 0x25, 0x15), RGBColor(0x5A, 0x4A, 0x1A))
tb = add_text_box(s, Inches(7.4), y + Inches(3.25), Inches(5.0), Inches(0.6))
set_text(tb.text_frame,
         "Trade-off: DeBERTa requires smaller batch sizes and slightly longer training time due to 26% more parameters. Worth it for the consistent quality gains.",
         size=11, color=LGRAY)

footer(s)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 15 — Conclusions
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); set_bg(s)
y = slide_title(s, "Conclusions", "Key Takeaways from Our ABSA Pipeline Evaluation")

conclusions = [
    ("Pipeline Works E2E", "Two-stage ATE → ASC extracts (aspect, sentiment) from raw text. In-domain: 84% ASC accuracy on restaurant, 80% on laptop."),
    ("ASC Transfers Well", "Sentiment classification is robust to domain shift. Cross-domain ASC accuracy: 78-92%, nearly matching in-domain."),
    ("ATE Is the Bottleneck", "Cross-domain ATE drops to 25-36% F1. ~85% of end-to-end errors originate from ATE failures."),
    ("Boundary Errors Dominate", "~37% of ATE errors are boundary errors (partial spans). Multi-word aspects are the hardest challenge."),
    ("DeBERTa > BERT", "DeBERTa outperforms BERT across all settings. Largest gain: +10.3 pp on cross-domain ATE. Worth the extra compute."),
    ("Neutral Class Is Hard", "Neutral F1: 0.51-0.64 across settings. Class imbalance (~20% Neutral) makes this the weakest class."),
]
for i, (title, text) in enumerate(conclusions):
    col = i % 2
    row = i // 2
    cx = Inches(0.7 + col * 6.25)
    cy = y + row * Inches(1.3)
    shape = add_rounded_rect(s, cx, cy, Inches(5.9), Inches(1.15), CARD_BG, CARD_BD)
    tb = add_text_box(s, cx + Inches(0.6), cy + Inches(0.08), Inches(5.1), Inches(1.0))
    tf = tb.text_frame
    set_text(tf, title, size=14, bold=True, color=WHITE)
    add_para(tf, text, size=11, color=GRAY)

    # Number circle
    shape = add_rounded_rect(s, cx + Inches(0.1), cy + Inches(0.15), Inches(0.38), Inches(0.38),
                             PURPLE, PURPLE)
    tb = add_text_box(s, cx + Inches(0.1), cy + Inches(0.17), Inches(0.38), Inches(0.35))
    set_text(tb.text_frame, str(i + 1), size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Bottom takeaway
shape = add_rounded_rect(s, Inches(0.7), y + Inches(4.1), Inches(11.9), Inches(0.7),
                         RGBColor(0x22, 0x1C, 0x35), RGBColor(0x45, 0x30, 0x70))
tb = add_text_box(s, Inches(0.9), y + Inches(4.15), Inches(11.5), Inches(0.6))
set_text(tb.text_frame,
         "Practical takeaway: For real-world ABSA deployment, invest in better ATE (domain-specific training data, joint modeling). ASC can be reused across domains with minimal degradation.",
         size=13, color=LGRAY)

footer(s)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 16 — Future Work
# ═══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); set_bg(s)
y = slide_title(s, "Future Work", "Directions for Improving End-to-End ABSA")

futures = [
    ("Joint ATE + ASC Modeling", "Train both tasks simultaneously with shared parameters to reduce error propagation.", "High Impact", GREEN),
    ("Larger Pre-trained Models", "DeBERTa-large, RoBERTa-large for richer representations. DeBERTa-base already gains +10pp.", "High Impact", GREEN),
    ("Data Augmentation", "Generate synthetic cross-domain training data. Back-translation and paraphrase generation.", "Medium Impact", YELLOW),
    ("Few-Shot Domain Adaptation", "Adapt to new domains with minimal labeled data. Meta-learning or prompt-tuning approaches.", "Medium Impact", YELLOW),
    ("Implicit Aspect Detection", "Handle aspects not explicitly mentioned. Requires reasoning beyond token spans.", "Exploratory", BLUE),
    ("LLM-Based ABSA", "Compare with GPT-4 / Claude zero-shot and few-shot. Can LLMs replace the two-stage pipeline?", "Exploratory", BLUE),
]
for i, (title, desc, impact, clr) in enumerate(futures):
    col = i % 2
    row = i // 2
    cx = Inches(0.7 + col * 6.25)
    cy = y + row * Inches(1.2)
    shape = add_rounded_rect(s, cx, cy, Inches(5.9), Inches(1.05), CARD_BG, CARD_BD)
    tb = add_text_box(s, cx + Inches(0.15), cy + Inches(0.08), Inches(5.6), Inches(0.9))
    tf = tb.text_frame
    set_text(tf, title, size=14, bold=True, color=WHITE)
    add_para(tf, desc, size=11, color=GRAY)
    p = add_para(tf, impact, size=10, bold=True, color=clr)

# Thank you
shape = add_rounded_rect(s, Inches(2.5), y + Inches(3.9), Inches(8.3), Inches(1.0),
                         RGBColor(0x22, 0x1C, 0x35), RGBColor(0x45, 0x30, 0x70))
tb = add_text_box(s, Inches(2.5), y + Inches(3.95), Inches(8.3), Inches(0.9))
tf = tb.text_frame
set_text(tf, "Thank You! Questions?", size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_para(tf, "Code and results available in the project repository", size=13, color=GRAY, align=PP_ALIGN.CENTER)

footer(s)

# ═══════════════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════════════
out_path = os.path.join(os.path.dirname(__file__), "ABSA_Final_Presentation.pptx")
prs.save(out_path)
print(f"Saved → {out_path}")
